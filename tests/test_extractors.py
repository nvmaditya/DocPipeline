from __future__ import annotations

import csv
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from docpipe.extractors import DocxExtractor, ExtractorRouter, HtmlExtractor, PdfExtractor, PptxExtractor, TextExtractor, XlsxExtractor


def test_pdf_extractor_reads_text_page(tmp_path: Path) -> None:
    pdf_path = tmp_path / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF extractor")
    doc.save(str(pdf_path))
    doc.close()

    extractor = PdfExtractor(scanned_threshold=5)
    records = extractor.extract(str(pdf_path))

    assert records
    assert "Hello PDF extractor" in records[0]["text"]
    assert records[0]["page_number"] == 1


def test_docx_extractor_reads_heading_paragraph_and_table(tmp_path: Path) -> None:
    docx_path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Section 1", level=1)
    doc.add_paragraph("This is a paragraph.")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    doc.save(str(docx_path))

    extractor = DocxExtractor()
    records = extractor.extract(str(docx_path))

    texts = [r["text"] for r in records]
    assert any("Section 1" in t for t in texts)
    assert any("This is a paragraph." in t for t in texts)
    assert any("A | B" in t for t in texts)


def test_pptx_extractor_reads_slide_text(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Revenue is up"
    prs.save(str(pptx_path))

    extractor = PptxExtractor()
    records = extractor.extract(str(pptx_path))

    assert records
    assert "Quarterly Review" in records[0]["text"]


def test_xlsx_and_csv_extractors(tmp_path: Path) -> None:
    xlsx_path = tmp_path / "data.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "SheetA"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    wb.save(str(xlsx_path))

    csv_path = tmp_path / "data.csv"
    with open(csv_path, "w", newline="", encoding="utf-8") as handle:
        writer = csv.writer(handle)
        writer.writerow(["City", "Temp"])
        writer.writerow(["Berlin", "12"])

    extractor = XlsxExtractor()
    xlsx_records = extractor.extract(str(xlsx_path))
    csv_records = extractor.extract(str(csv_path))

    assert any("Alice" in r["text"] for r in xlsx_records)
    assert any("Berlin" in r["text"] for r in csv_records)


def test_html_and_text_extractors(tmp_path: Path) -> None:
    html_path = tmp_path / "sample.html"
    html_path.write_text("<html><body><h1>Title</h1><p>Body text.</p></body></html>", encoding="utf-8")

    txt_path = tmp_path / "sample.txt"
    txt_path.write_text("plain text content", encoding="utf-8")

    html_records = HtmlExtractor().extract(str(html_path))
    txt_records = TextExtractor().extract(str(txt_path))

    assert "Title" in html_records[0]["text"]
    assert "plain text content" in txt_records[0]["text"]


def test_router_supports_phase2_extensions(tmp_path: Path) -> None:
    router = ExtractorRouter()
    assert isinstance(router.route(str(tmp_path / "a.pdf")), PdfExtractor)
    assert isinstance(router.route(str(tmp_path / "a.docx")), DocxExtractor)
    assert isinstance(router.route(str(tmp_path / "a.pptx")), PptxExtractor)
    assert isinstance(router.route(str(tmp_path / "a.xlsx")), XlsxExtractor)
    assert isinstance(router.route(str(tmp_path / "a.csv")), XlsxExtractor)
    assert isinstance(router.route(str(tmp_path / "a.html")), HtmlExtractor)
    assert isinstance(router.route(str(tmp_path / "a.txt")), TextExtractor)
