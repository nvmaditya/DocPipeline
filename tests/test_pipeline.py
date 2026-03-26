from __future__ import annotations

import csv
from pathlib import Path

import fitz
import numpy as np
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import pytest

import docpipe.pipeline as pipeline_mod


def _make_pdf(path: Path, text: str) -> None:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    doc.save(str(path))
    doc.close()


def _make_docx(path: Path, heading: str, text: str) -> None:
    doc = Document()
    doc.add_heading(heading, level=1)
    doc.add_paragraph(text)
    doc.save(str(path))


def _make_pptx(path: Path, title: str, body: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    prs.save(str(path))


def _make_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Metric", "Value"])
    ws.append(["Growth", "12%"])
    wb.save(str(path))


def _make_csv(path: Path) -> None:
    with open(path, "w", newline="", encoding="utf-8") as handle:
        writer = csv.writer(handle)
        writer.writerow(["Region", "Sales"])
        writer.writerow(["EMEA", "120"])


def test_pipeline_ingest_and_query(tmp_path: Path) -> None:
    docs_dir = tmp_path / "docs"
    docs_dir.mkdir()

    _make_pdf(docs_dir / "sample.pdf", "Revenue increased in Q3.")
    _make_docx(docs_dir / "sample.docx", "Summary", "Revenue increased and costs dropped.")
    _make_pptx(docs_dir / "sample.pptx", "Results", "Revenue in presentation")
    _make_xlsx(docs_dir / "sample.xlsx")
    _make_csv(docs_dir / "sample.csv")
    (docs_dir / "sample.html").write_text("<html><body><p>Revenue in html</p></body></html>", encoding="utf-8")
    (docs_dir / "sample.txt").write_text("Revenue in text file", encoding="utf-8")

    store_dir = tmp_path / "store"
    faiss_path = (store_dir / "faiss.index").as_posix()
    sqlite_path = (store_dir / "metadata.db").as_posix()
    cfg_path = tmp_path / "config.yaml"
    cfg_path.write_text(
        "\n".join(
            [
                "extraction:",
                "  scanned_threshold: 10",
                "  ocr_engine: none",
                "  ocr_language: eng",
                "chunking:",
                "  strategy: semantic",
                "  chunk_size: 256",
                "  chunk_overlap: 32",
                "  semantic_threshold: 0.5",
                "embedding:",
                "  model: sentence-transformers/all-MiniLM-L6-v2",
                "  batch_size: 8",
                "  device: cpu",
                "  normalize: true",
                "faiss:",
                "  index_type: hnsw",
                "  hnsw_m: 16",
                "  hnsw_ef_construction: 80",
                "  hnsw_ef_search: 32",
                "store:",
                f"  faiss_path: {faiss_path}",
                f"  sqlite_path: {sqlite_path}",
                "query:",
                "  top_k: 5",
                "  score_threshold: 0.0",
            ]
        ),
        encoding="utf-8",
    )

    class DummyEmbedder:
        def __init__(self, *args, **kwargs) -> None:
            pass

        def encode(self, texts):
            vectors = []
            for text in texts:
                length = float(len(text))
                alpha = float(sum(1 for c in text.lower() if c.isalpha()))
                digits = float(sum(1 for c in text if c.isdigit()))
                vectors.append([length, alpha, digits, 1.0])
            arr = np.asarray(vectors, dtype=np.float32)
            norms = np.linalg.norm(arr, axis=1, keepdims=True)
            norms[norms == 0] = 1.0
            return arr / norms

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(pipeline_mod, "Embedder", DummyEmbedder)

    pipe = pipeline_mod.Pipeline(config=str(cfg_path))
    try:
        ingested = pipe.ingest(str(docs_dir))
        assert ingested >= 2
    finally:
        pipe.close()

    pipe2 = pipeline_mod.Pipeline(config=str(cfg_path))
    try:
        results = pipe2.search("revenue", top_k=5)
        assert results
        assert any("Revenue" in r["chunk_text"] or "revenue" in r["chunk_text"] for r in results)

        stats = pipe2.stats()
        assert stats["documents"] == ingested
        assert stats["chunks"] >= 2
    finally:
        pipe2.close()
        monkeypatch.undo()
