from __future__ import annotations

from typing import Any, Dict, List

from pptx import Presentation

from .base import BaseExtractor


class PptxExtractor(BaseExtractor):
    def extract(self, file_path: str) -> List[Dict[str, Any]]:
        prs = Presentation(file_path)
        records: List[Dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = (shape.text or "").strip()
                    if text:
                        parts.append(text)
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                parts.append(notes_text)
            if parts:
                records.append(
                    {
                        "text": "\n".join(parts),
                        "page_number": idx,
                        "heading_context": f"Slide {idx}",
                    }
                )
        return records
